from tkinter import *
from PIL import ImageTk, Image
from tkinter.filedialog import askopenfile
import pyttsx3
import PyPDF2
from pptx import Presentation
import threading
import sys

def playThroughThread():
    global thread
    thread = threading.Thread(target=play, args=())
    thread.daemon = True  # Daemonize thread
    thread.start()  # Start the execution

def play():
    message='Message: Reading '+file_name_display+' pages from: '+start_pg_entry.get()+' to: '+end_pg_entry.get()
    message_entry.delete(0, END)
    message_entry.insert(0,message)
    document = open(file_name, 'rb')
    speaker = pyttsx3.init()
    text_runs = []
    if var.get() == 1:
        ppt_reader = Presentation(document)
        print(len(ppt_reader.slides))
        for slide in ppt_reader.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
    elif var.get()==2:
        pdf_reader = PyPDF2.PdfFileReader(document)
        pages = pdf_reader.numPages
        for page_no in range(pages):
            single_page = pdf_reader.getPage(page_no)
            text_runs.append(single_page.extractText().replace('\n',' '))

    document.close()
    speaker.say(text_runs)
    speaker.runAndWait()



def stop():
    message='Message: Stopped player'
    message_entry.delete(0, END)
    message_entry.insert(0,message)
    sys.exit()


def open_file():
    global file_name, file_name_display
    if var.get() == 1:  # PPT
        file = askopenfile(mode='r', filetypes=[('Python Files', '*.pptx')])
        ppt_document = open(file.name, 'rb')
        ppt_reader = Presentation(ppt_document)
        start_pg_entry.config(state='normal')
        start_pg_entry.delete(0,END)
        start_pg_entry.insert(0,1)
        end_pg_entry.config(state='normal')
        end_pg_entry.delete(0,END)
        end_pg_entry.insert(0,len(ppt_reader.slides))
        ppt_document.close()

    elif var.get() == 2:  # PDF
        file = askopenfile(mode='r', filetypes=[('Python Files', '*.pdf')])
        pdf_document = open(file.name, 'rb')
        pdf_reader = PyPDF2.PdfFileReader(pdf_document)
        start_pg_entry.config(state='normal')
        start_pg_entry.delete(0, END)
        start_pg_entry.insert(0, 1)
        end_pg_entry.config(state='normal')
        end_pg_entry.delete(0, END)
        end_pg_entry.insert(0, pdf_reader.numPages)
        pdf_document.close()

    if file is not None:
        file_name = file.name
        print(file.name)
        file_name_display=file_name[file_name.rfind('/')+1: len(file_name)]
        message = 'Message: Choosen file: ' + file_name_display
    else:
        message = 'Message: Error occured while chooing file'
    message_entry.delete(0, END)
    message_entry.insert(0, message)




root = Tk()
root.title('Read out loud')

img = ImageTk.PhotoImage(Image.open('image.jpg'))
img_lbl = Label(image=img)

var = IntVar()
#Row1 radio buttons
R1 = Radiobutton(root, text="PPT", variable=var, value=1)
R1.select()
R2 = Radiobutton(root, text="PDF", variable=var, value=2)

open_file_btn = Button(root, text='Open File', command=lambda: open_file(), borderwidth=5, fg="white", bg="purple", width=52)

start_lbl = Label(root, text='Start Page: ')
start_pg_entry= Entry(root, width=10, borderwidth=5, state='readonly')


end_lbl=Label(root, text='End Page: ')
end_pg_entry= Entry(root, width=10, borderwidth=5, state='readonly')

message = 'Message:'
message_entry = Entry(root, fg="white", bg="grey", width=60)
message_entry.delete(0, END)
message_entry.insert(0, message)


page_selection_btn = Button(root, text="▶", command=playThroughThread, borderwidth=5, fg="white", bg="green", width=25,height=5 )
stop_btn = Button(root, text="⏹", command=stop, borderwidth=2, fg="white", bg="red", width=25, height=5)


#packing
img_lbl.grid(row=0, column=0, columnspan=4)
R1.grid(row=1, column=1 )
R2.grid(row=1, column=2)
open_file_btn.grid(row=2, column=0, padx=5, pady=5, columnspan=4)
start_lbl.grid(row=3,column=0)
start_pg_entry.grid(row=3,column=1)
end_lbl.grid(row=3,column=2)
end_pg_entry.grid(row=3,column=3)
page_selection_btn.grid(row=4, column=0,  padx=2, columnspan=2,pady=5)
stop_btn.grid(row=4, column=2, padx=2, columnspan=2,pady=5)
message_entry.grid(row=10, column=0, columnspan=4, pady=4)


root.mainloop()





